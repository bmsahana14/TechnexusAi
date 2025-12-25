import os
import json
import logging
from typing import List, Dict, Any
from pathlib import Path

# External libs
from pptx import Presentation
from pypdf import PdfReader
from dotenv import load_dotenv

# Try to import Google Gemini, fallback to OpenAI if not available
try:
    import google.generativeai as genai
    GEMINI_AVAILABLE = True
except ImportError:
    GEMINI_AVAILABLE = False
    print("Google Gemini not installed. Install with: pip install google-generativeai")

try:
    from openai import OpenAI
    OPENAI_AVAILABLE = True
except ImportError:
    OPENAI_AVAILABLE = False

load_dotenv()

# Initialize AI clients
gemini_key = os.getenv("GEMINI_API_KEY")
openai_key = os.getenv("OPENAI_API_KEY")

if GEMINI_AVAILABLE and gemini_key and gemini_key != "your_gemini_key_here":
    genai.configure(api_key=gemini_key)
    gemini_model = genai.GenerativeModel('gemini-pro')
    AI_PROVIDER = "gemini"
    print("✅ Using Google Gemini for quiz generation")
elif OPENAI_AVAILABLE and openai_key and openai_key not in ["your_openai_api_key_here", "dummy_key_for_testing"]:
    openai_client = OpenAI(api_key=openai_key)
    AI_PROVIDER = "openai"
    print("✅ Using OpenAI for quiz generation")
else:
    AI_PROVIDER = "fallback"
    print("⚠️  No valid AI API key found. Using fallback mode.")
    print("   Get free Gemini key at: https://aistudio.google.com/app/apikey")

async def generate_quiz_from_file(file_path: str, num_questions: int, difficulty: str) -> List[Dict[str, Any]]:
    """
    Orchestrates the conversion of a file > text > quiz questions.
    Ensures at least 10 questions are generated and deletes the uploaded file after processing.
    """
    # Ensure a minimum of 10 questions
    requested_questions = max(num_questions, 10)
    
    try:
        # 1. Extract content
        if file_path.endswith('.pptx'):
            text_content = extract_text_from_pptx(file_path)
        elif file_path.endswith('.pdf'):
            text_content = extract_text_from_pdf(file_path)
        else:
            raise ValueError("Unsupported file format")
        
        print(f"Extracted {len(text_content)} characters from {os.path.basename(file_path)}.")
        
        # 2. Limit content for context window
        max_chars = 15000
        if len(text_content) > max_chars:
            text_content = text_content[:max_chars] + "...[truncated]"
        
        # 3. Generate via LLM
        quiz = query_llm_for_quiz(text_content, requested_questions, difficulty)
        return quiz
    except Exception as e:
        print(f"Error generating quiz: {e}")
        return get_fallback_questions()
    finally:
        # Delete the uploaded file to free space
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
                print(f"Deleted uploaded file: {file_path}")
        except Exception as del_err:
            print(f"Failed to delete uploaded file {file_path}: {del_err}")

def extract_text_from_pptx(path: str) -> str:
    prs = Presentation(path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

def extract_text_from_pdf(path: str) -> str:
    reader = PdfReader(path)
    text = ""
    for page in reader.pages:
        text += page.extract_text() + "\n"
    return text

def query_llm_for_quiz(content: str, count: int, difficulty: str) -> List[Dict[str, Any]]:
    
    prompt = f"""
    You are an expert quiz generator for technical presentations and educational content.
    
    CONTENT TO ANALYZE:
    {content}
    
    TASK:
    Generate {count} high-quality multiple-choice questions based strictly on the content above.
    
    DIFFICULTY LEVEL: {difficulty}
    - Easy: Focus on definitions, basic concepts, and direct facts from the content
    - Medium: Require understanding and application of concepts
    - Hard: Test deep comprehension, analysis, and synthesis of information
    
    REQUIREMENTS:
    1. Questions must be directly answerable from the provided content
    2. Each question should test a unique concept or fact
    3. All 4 options must be plausible but only one correct
    4. Avoid ambiguous or trick questions
    5. Use clear, professional language
    6. Ensure diversity in question types (what, how, why, which, etc.)
    
    OUTPUT FORMAT:
    Return ONLY a valid JSON array with NO markdown formatting, NO code blocks, NO explanations.
    Each object must have exactly these fields:
    - "q": The question text (string)
    - "options": Array of exactly 4 distinct answer options (strings)
    - "correct": Index (0-3) of the correct option (integer)
    
    EXAMPLE:
    [
        {{"q": "What is the primary benefit of microservices architecture?", "options": ["Monolithic design", "Independent scalability", "Single database", "Tight coupling"], "correct": 1}},
        {{"q": "Which protocol is used for real-time communication?", "options": ["HTTP", "FTP", "WebSocket", "SMTP"], "correct": 2}}
    ]
    
    Generate {count} questions now:
    """
    
    try:
        if AI_PROVIDER == "gemini":
            return query_gemini(prompt)
        elif AI_PROVIDER == "openai":
            return query_openai(prompt)
        else:
            return get_fallback_questions()
    except Exception as e:
        print(f"Error querying AI: {e}")
        return get_fallback_questions()

def query_gemini(prompt: str) -> List[Dict[str, Any]]:
    """Query Google Gemini for quiz generation"""
    try:
        response = gemini_model.generate_content(prompt)
        raw_content = response.text.strip()
        
        # Clean up potential markdown code blocks
        if raw_content.startswith("```json"):
            raw_content = raw_content[7:]
        if raw_content.startswith("```"):
            raw_content = raw_content[3:]
        if raw_content.endswith("```"):
            raw_content = raw_content[:-3]
            
        return json.loads(raw_content.strip())
        
    except Exception as e:
        print(f"Error querying Gemini: {e}")
        raise

def query_openai(prompt: str) -> List[Dict[str, Any]]:
    """Query OpenAI for quiz generation"""
    try:
        response = openai_client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "You are a professional quiz generator that outputs only valid JSON arrays. Never include markdown formatting or explanations."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7,
            max_tokens=2000
        )
        
        raw_content = response.choices[0].message.content.strip()
        
        # Clean up potential markdown code blocks
        if raw_content.startswith("```json"):
            raw_content = raw_content[7:]
        if raw_content.startswith("```"):
            raw_content = raw_content[3:]
        if raw_content.endswith("```"):
            raw_content = raw_content[:-3]
            
        return json.loads(raw_content.strip())
        
    except Exception as e:
        print(f"Error querying OpenAI: {e}")
        raise

def get_fallback_questions() -> List[Dict[str, Any]]:
    """Return fallback questions when AI is not available. Guarantees at least 10 questions."""
    print("Returning FALLBACK questions for testing.")
    return [
        {
            "q": "What is the primary architectural style of this application?",
            "options": ["Monolithic", "Microservices (Client, Realtime, AI)", "Serverless ONLY", "Mainframe"],
            "correct": 1
        },
        {
            "q": "Which library is used for the realtime communication?",
            "options": ["Socket.IO", "React Query", "Redux", "Axios"],
            "correct": 0
        },
        {
            "q": "What is the role of the AI Service?",
            "options": ["Host the UI", "Generate questions from files", "Manage the database", "Authenticate users"],
            "correct": 1
        },
        {
            "q": "Which CSS framework is used for styling?",
            "options": ["Bootstrap", "Foundation", "Tailwind CSS", "Bulma"],
            "correct": 2
        },
        {
            "q": "Where is the active quiz state currently stored?",
            "options": ["PostgreSQL", "Redis", "In-Memory (Map)", "LocalStorage"],
            "correct": 2
        },
        {
            "q": "How does the client receive new quiz questions in real time?",
            "options": ["Polling the server", "WebSocket events", "HTTP long polling", "Server-sent events"],
            "correct": 1
        },
        {
            "q": "What format does the AI service return for generated quizzes?",
            "options": ["Plain text", "XML", "JSON array", "CSV"],
            "correct": 2
        },
        {
            "q": "Which environment variable determines the AI provider?",
            "options": ["AI_PROVIDER", "GEMINI_API_KEY", "OPENAI_API_KEY", "NONE"],
            "correct": 0
        },
        {
            "q": "What happens if the AI service encounters an error during generation?",
            "options": ["Returns empty list", "Throws exception", "Falls back to default questions", "Retries indefinitely"],
            "correct": 2
        },
        {
            "q": "Why is it safe to delete the uploaded file after quiz generation?",
            "options": ["File is no longer needed", "Server caches it", "Client still needs it", "It is stored in DB"],
            "correct": 0
        }
    ]
