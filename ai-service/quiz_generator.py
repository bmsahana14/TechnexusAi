import os
import json
import logging
from typing import List, Dict, Any
from pathlib import Path

# External libs
from pptx import Presentation
from pypdf import PdfReader
from dotenv import load_dotenv

# Try to import Google Gemini, fallback to OpenAI if not available
try:
    import google.generativeai as genai
    GEMINI_AVAILABLE = True
except ImportError:
    GEMINI_AVAILABLE = False
    print("Google Gemini not installed. Install with: pip install google-generativeai")

try:
    from openai import OpenAI
    OPENAI_AVAILABLE = True
except ImportError:
    OPENAI_AVAILABLE = False

load_dotenv()

# Initialize AI clients
# Initialize AI clients
gemini_key = os.getenv("GEMINI_API_KEY")
openai_key = os.getenv("OPENAI_API_KEY")

print("=" * 60)
print("AI SERVICE INITIALIZATION")
print("=" * 60)

# DEBUG LOGGING TO FILE
with open("ai_init_log.txt", "w") as f:
    f.write("Starting initialization...\n")
    f.write(f"Gemini Key Present: {bool(gemini_key)}\n")
    if gemini_key:
        f.write(f"Gemini Key Start: {gemini_key[:5]}...\n")

if GEMINI_AVAILABLE and gemini_key and gemini_key != "your_gemini_key_here":
    try:
        # Explicitly configure with transport='rest' to avoid GRPC issues
        genai.configure(api_key=gemini_key, transport='rest')
        
        # Use simple model name
        gemini_model = genai.GenerativeModel('models/gemini-flash-latest')
        
        # Test generation to confirm it works
        # response = gemini_model.generate_content("hello")
        
        AI_PROVIDER = "gemini"
        print("OK: Using Google Gemini Flash (Latest) for quiz generation")
        print(f"   API Key: {gemini_key[:10]}...{gemini_key[-4:]}")
        
        with open("ai_init_log.txt", "a") as f:
            f.write("SUCCESS: Gemini initialized!\n")
            
    except Exception as e:
        print(f"ERROR: Gemini initialization failed: {e}")
        with open("ai_init_log.txt", "a") as f:
            f.write(f"ERROR: Initialization failed: {e}\n")
            import traceback
            traceback.print_exc(file=f)
            
        AI_PROVIDER = "fallback"

elif OPENAI_AVAILABLE and openai_key and openai_key not in ["your_openai_api_key_here", "dummy_key_for_testing"]:

    try:
        openai_client = OpenAI(api_key=openai_key)
        AI_PROVIDER = "openai"
        print("OK: Using OpenAI for quiz generation")
        print(f"   API Key: {openai_key[:10]}...{openai_key[-4:]}")
    except Exception as e:
        print(f"ERROR: OpenAI initialization failed: {e}")
        AI_PROVIDER = "fallback"
else:
    AI_PROVIDER = "fallback"
    print("WARNING: No valid AI API key found. Using fallback mode.")
    print("   Get free Gemini key at: https://aistudio.google.com/app/apikey")
    if not GEMINI_AVAILABLE:
        print("   Note: google-generativeai not installed")
    if gemini_key:
        print(f"   Gemini key status: {gemini_key[:20]}...")

print("=" * 60)


async def generate_quiz_from_file(file_path: str, num_questions: int, difficulty: str) -> List[Dict[str, Any]]:
    """
    Orchestrates the conversion of a file > text > quiz questions
    """
    
    print(f"\n{'='*60}")
    print(f"QUIZ GENERATION REQUEST")
    print(f"{'='*60}")
    print(f"File: {file_path}")
    print(f"Questions: {num_questions}")
    print(f"Difficulty: {difficulty}")
    print(f"AI Provider: {AI_PROVIDER}")
    
    # 1. Extract content
    try:
        if file_path.endswith('.pptx'):
            print("Extracting text from PPTX...")
            text_content = extract_text_from_pptx(file_path)
        elif file_path.endswith('.pdf'):
            print("Extracting text from PDF...")
            text_content = extract_text_from_pdf(file_path)
        else:
            raise ValueError("Unsupported file format")
            
        print(f"OK: Extracted {len(text_content)} characters from file")
        
        if len(text_content) < 100:
            print(f"WARNING: Very little text extracted ({len(text_content)} chars)")
            print(f"   Preview: {text_content[:200]}")
        else:
            print(f"   Preview: {text_content[:200]}...")
        
        # 2. Limit content for context window
        max_chars = 15000 
        if len(text_content) > max_chars:
            print(f"📝 Truncating content from {len(text_content)} to {max_chars} characters")
            text_content = text_content[:max_chars] + "...[truncated]"

        # 3. Generate via LLM
        print(f"Sending to {AI_PROVIDER.upper()} for quiz generation...")
        questions = query_llm_for_quiz(text_content, num_questions, difficulty)
        
        print("OK: Generated questions successfully")
        print(f"{'='*60}\n")
        return questions
        
    except Exception as e:
        import traceback
        print("ERROR in quiz generation:")
        traceback.print_exc()
        print(f"Error: {e}")
        print("WARNING: Returning fallback questions")
        print(f"{'='*60}\n")
        return get_fallback_questions()

def extract_text_from_pptx(path: str) -> str:
    prs = Presentation(path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

def extract_text_from_pdf(path: str) -> str:
    reader = PdfReader(path)
    text = ""
    for page in reader.pages:
        text += page.extract_text() + "\n"
    return text

def query_llm_for_quiz(content: str, count: int, difficulty: str) -> List[Dict[str, Any]]:
    
    prompt = f"""
    You are an expert quiz generator for technical presentations and educational content.
    
    CONTENT TO ANALYZE:
    {content}
    
    TASK:
    Generate {count} high-quality multiple-choice questions based strictly on the content above.
    
    DIFFICULTY LEVEL: {difficulty}
    - Easy: Focus on definitions, basic concepts, and direct facts from the content
    - Medium: Require understanding and application of concepts
    - Hard: Test deep comprehension, analysis, and synthesis of information
    
    REQUIREMENTS:
    1. Questions must be directly answerable from the provided content
    2. Each question should test a unique concept or fact
    3. All 4 options must be plausible but only one correct
    4. Avoid ambiguous or trick questions
    5. Use clear, professional language
    6. Ensure diversity in question types (what, how, why, which, etc.)
    
    OUTPUT FORMAT:
    Return ONLY a valid JSON array with NO markdown formatting, NO code blocks, NO explanations.
    Each object must have exactly these fields:
    - "q": The question text (string)
    - "options": Array of exactly 4 distinct answer options (strings)
    - "correct": Index (0-3) of the correct option (integer)
    
    EXAMPLE:
    [
        {{"q": "What is the primary benefit of microservices architecture?", "options": ["Monolithic design", "Independent scalability", "Single database", "Tight coupling"], "correct": 1}},
        {{"q": "Which protocol is used for real-time communication?", "options": ["HTTP", "FTP", "WebSocket", "SMTP"], "correct": 2}}
    ]
    
    Generate {count} questions now:
    """
    
    try:
        if AI_PROVIDER == "gemini":
            return query_gemini(prompt)
        elif AI_PROVIDER == "openai":
            return query_openai(prompt)
        else:
            return get_fallback_questions()
    except Exception as e:
        print(f"ERROR: AI Generation failed: {e}")
        print("WARNING: Falling back to generic questions!")
        return get_fallback_questions()


def query_gemini(prompt: str) -> List[Dict[str, Any]]:
    """Query Google Gemini for quiz generation"""
    try:
        response = gemini_model.generate_content(prompt)
        raw_content = response.text.strip()
        print(f"DEBUG: Raw Content from Gemini: {raw_content[:200]}...")

        
        # Clean up potential markdown code blocks
        if raw_content.startswith("```json"):
            raw_content = raw_content[7:]
        if raw_content.startswith("```"):
            if "json" in raw_content[:10]:
                 raw_content = raw_content[raw_content.find("["):]
            else:
                 raw_content = raw_content[3:]
        if raw_content.endswith("```"):
            raw_content = raw_content[:raw_content.rfind("]")+1]
            
        # Robust JSON extraction
        start_idx = raw_content.find("[")
        end_idx = raw_content.rfind("]")
        if start_idx != -1 and end_idx != -1:
            raw_content = raw_content[start_idx:end_idx+1]

        return json.loads(raw_content.strip())
        
    except Exception as e:
        print(f"Error querying Gemini: {e}")
        raise

def query_openai(prompt: str) -> List[Dict[str, Any]]:
    """Query OpenAI for quiz generation"""
    try:
        response = openai_client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "You are a professional quiz generator that outputs only valid JSON arrays. Never include markdown formatting or explanations."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7,
            max_tokens=2000
        )
        
        raw_content = response.choices[0].message.content.strip()
        
        # Clean up potential markdown code blocks
        if raw_content.startswith("```json"):
            raw_content = raw_content[7:]
        if raw_content.startswith("```"):
            raw_content = raw_content[3:]
        if raw_content.endswith("```"):
            raw_content = raw_content[:-3]
            
        return json.loads(raw_content.strip())
        
    except Exception as e:
        print(f"Error querying OpenAI: {e}")
        raise

def get_fallback_questions() -> List[Dict[str, Any]]:
    """Return fallback questions when AI is not available"""
    print("Returning FALLBACK questions for testing.")
    return [
        {
            "q": "What is the primary architectural style of this application?",
            "options": ["Monolithic", "Microservices (Client, Realtime, AI)", "Serverless ONLY", "Mainframe"],
            "correct": 1
        },
        {
            "q": "Which library is used for the realtime communication?",
            "options": ["Socket.IO", "React Query", "Redux", "Axios"],
            "correct": 0
        },
        {
            "q": "What is the role of the AI Service?",
            "options": ["Host the UI", "Generate questions from files", "Manage the database", "Authenticate users"],
            "correct": 1
        },
        {
            "q": "Which CSS framework is used for styling?",
            "options": ["Bootstrap", "Foundation", "Tailwind CSS", "Bulma"],
            "correct": 2
        },
        {
            "q": "Where is the active quiz state currently stored?",
            "options": ["PostgreSQL", "Redis", "In-Memory (Map)", "LocalStorage"],
            "correct": 2
        }
    ]
